"""Presentation rendering/orchestration for weekly status PPT output."""

from __future__ import annotations

from datetime import date, timedelta
from pathlib import Path
from typing import Callable

from pptx import Presentation

from .ai_priority import format_ai_priority, get_ai_priority
from .data_service import get_feature_items, get_user_stories_from_rally_api
from .milestone_service import get_milestone_release_date, get_milestone_release_mapping
from .config import debug
from .story_utils import extract_current_status, extract_primary_milestone_id

DEFAULT_TEMPLATE = Path("templates") / "Weekly-Status-Template.pptx"
DEFAULT_OUTPUT = Path("generated_ppts") / "Weekly-Status-Template-Filled.pptx"


def _emit_progress(progress_callback: Callable[[str], None] | None, message: str) -> None:
    if not progress_callback:
        return
    try:
        progress_callback(message)
    except Exception:
        # Progress updates should never break PPT generation.
        pass


def week_range(ref_date: date | None = None) -> tuple[date, date]:
    if ref_date is None:
        ref_date = date.today()
    monday = ref_date - timedelta(days=ref_date.weekday())
    friday = monday + timedelta(days=4)
    return monday, friday


def format_range(monday: date, friday: date, date_format: str) -> tuple[str, str]:
    return monday.strftime(date_format), friday.strftime(date_format)


def replace_text_in_shape(shape, from_text: str, to_text: str, team_name: str = "") -> bool:
    if not shape.has_text_frame:
        return False

    text_frame = shape.text_frame
    replaced = False

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run_text = run.text
            if not run_text:
                continue
            if "{FROM_DATE}" in run_text or "{TO_DATE}" in run_text or "{TEAM}" in run_text:
                run.text = run_text.replace("{FROM_DATE}", from_text).replace("{TO_DATE}", to_text).replace("{TEAM}", team_name)
                replaced = True

    if replaced:
        return True

    import re

    pattern = re.compile(r"(Project\s+Status\s*[–-]\s*From)\s+To", re.IGNORECASE)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run_text = run.text
            if run_text and pattern.search(run_text):
                run.text = pattern.sub(rf"\1 {from_text} To {to_text}", run_text)
                return True
    return False


def apply_team_text_to_slide(slide, from_text: str, to_text: str, team_name: str) -> None:
    """Apply date placeholders and normalize stamped team names on a slide.

    Why: some sections are cloned from slides that were already rendered for a
    different team, so we must replace both placeholders and existing team text.
    """
    for shape in slide.shapes:
        replace_text_in_shape(shape, from_text, to_text, team_name)
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "{TEAM}" in run.text:
                    run.text = run.text.replace("{TEAM}", team_name)
                if team_name != "Eagle" and "Eagle" in run.text:
                    run.text = run.text.replace("Eagle", team_name)
                if team_name != "Panther" and "Panther" in run.text:
                    run.text = run.text.replace("Panther", team_name)


def apply_run_style(target_run, source_run) -> None:
    if source_run is None or target_run is None or source_run.font is None:
        return
    if source_run.font.name:
        target_run.font.name = source_run.font.name
    if source_run.font.size:
        target_run.font.size = source_run.font.size
    if source_run.font.bold is not None:
        target_run.font.bold = source_run.font.bold
    if source_run.font.italic is not None:
        target_run.font.italic = source_run.font.italic
    if source_run.font.color:
        try:
            if source_run.font.color.rgb is not None:
                target_run.font.color.rgb = source_run.font.color.rgb
        except Exception:
            try:
                if source_run.font.color.theme_color is not None:
                    target_run.font.color.theme_color = source_run.font.color.theme_color
                    if source_run.font.color.brightness is not None:
                        target_run.font.color.brightness = source_run.font.color.brightness
            except Exception:
                pass


def set_cell_text(cell, text: str, source_run, source_paragraph=None) -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    if source_paragraph is not None and source_paragraph.alignment is not None:
        paragraph.alignment = source_paragraph.alignment
    run = paragraph.add_run()
    run.text = text
    apply_run_style(run, source_run)


def fill_user_story_table(
    presentation: Presentation,
    user_stories: list[dict],
    slide_idx: int = 1,
    global_counter_start: int = 1,
    upcoming_release: str | None = None,
) -> tuple[int, int]:
    if len(presentation.slides) <= slide_idx:
        return 0, global_counter_start

    slide = presentation.slides[slide_idx]
    table_shapes = [shape for shape in slide.shapes if shape.has_table]
    if not table_shapes:
        return 0, global_counter_start

    table = table_shapes[0].table
    if len(table.rows) < 3:
        return 0, global_counter_start

    sample_row = table.rows[1]
    sample_runs = []
    sample_paragraphs = []
    num_cols = len(table.columns)
    for col_idx in range(min(num_cols, 6)):
        sample_cell = sample_row.cells[col_idx]
        sample_run = None
        sample_paragraph = None
        if sample_cell.text_frame.paragraphs and sample_cell.text_frame.paragraphs[0].runs:
            sample_paragraph = sample_cell.text_frame.paragraphs[0]
            sample_run = sample_paragraph.runs[0]
        sample_runs.append(sample_run)
        sample_paragraphs.append(sample_paragraph)

    max_rows = len(table.rows) - 2
    filled = 0
    global_counter = global_counter_start

    for idx, item in enumerate(user_stories[:max_rows], start=1):
        row = table.rows[idx + 1]
        formatted_id = str(item.get("FormattedID") or item.get("FormattedId") or item.get("ID") or item.get("id") or "")
        name = item.get("Name") or item.get("name") or ""

        business_value = (
            item.get("c_BusinessValue")
            or item.get("c_businessValue")
            or item.get("BusinessValue")
            or item.get("businessValue")
            or item.get("Business Value")
        )

        if business_value is not None and str(business_value).strip() not in ("", "None", "null"):
            priority = str(business_value)
        elif name:
            priority = get_ai_priority(name)
        else:
            priority = format_ai_priority("2")

        current_status = extract_current_status(item)
        milestone = ""
        milestone_id = extract_primary_milestone_id(item)

        if milestone_id:
            mapped = get_milestone_release_date(milestone_id)
            if mapped:
                milestone = mapped
            elif upcoming_release:
                milestone = upcoming_release
            else:
                milestone = milestone_id

        set_cell_text(row.cells[0], str(global_counter), sample_runs[0], sample_paragraphs[0])
        set_cell_text(row.cells[1], str(formatted_id), sample_runs[1], sample_paragraphs[1])
        set_cell_text(row.cells[2], str(name), sample_runs[2], sample_paragraphs[2])

        if num_cols > 3:
            set_cell_text(row.cells[3], str(priority), sample_runs[3] if len(sample_runs) > 3 else sample_runs[0], sample_paragraphs[3] if len(sample_paragraphs) > 3 else sample_paragraphs[0])
        if num_cols > 4:
            set_cell_text(row.cells[4], str(current_status), sample_runs[4] if len(sample_runs) > 4 else sample_runs[0], sample_paragraphs[4] if len(sample_paragraphs) > 4 else sample_paragraphs[0])
        if num_cols > 5:
            set_cell_text(row.cells[5], str(milestone), sample_runs[5] if len(sample_runs) > 5 else sample_runs[0], sample_paragraphs[5] if len(sample_paragraphs) > 5 else sample_paragraphs[0])

        filled += 1
        global_counter += 1

    delete_example_row(table)
    delete_unused_rows(table, filled)
    return filled, global_counter


def delete_example_row(table) -> None:
    if len(table.rows) < 2:
        return
    tr = table.rows[1]._tr
    tr.getparent().remove(tr)


def delete_unused_rows(table, filled_rows: int) -> None:
    if len(table.rows) <= 1:
        return
    target_rows = 1 + max(filled_rows, 0)
    while len(table.rows) > target_rows:
        tr = table.rows[len(table.rows) - 1]._tr
        tr.getparent().remove(tr)


def clone_slide(presentation: Presentation, source_idx: int, insert_after_idx: int) -> int:
    from lxml import etree

    source_slide = presentation.slides[source_idx]
    new_slide = presentation.slides.add_slide(source_slide.slide_layout)

    for shape in source_slide.shapes:
        if shape.is_placeholder:
            continue
        el = etree.fromstring(etree.tostring(shape.element))
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    xml_slides = presentation.slides._sldIdLst
    new_slide_el = xml_slides[-1]
    xml_slides.remove(new_slide_el)
    xml_slides.insert(insert_after_idx + 1, new_slide_el)
    return insert_after_idx + 1


def fill_project_status_dates(
    template_path: Path,
    output_path: Path,
    ref_date: date | None = None,
    date_format: str = "%b %d, %Y",
    team: str = "",
) -> Path:
    monday, friday = week_range(ref_date)
    from_text, to_text = format_range(monday, friday, date_format)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(template_path))
    replaced = False
    for slide in presentation.slides:
        for shape in slide.shapes:
            replaced = replace_text_in_shape(shape, from_text, to_text, team) or replaced

    if not replaced:
        raise RuntimeError(
            "Could not find 'Project Status – From To' text or placeholders to replace. "
            "Add {FROM_DATE}/{TO_DATE}/{TEAM} or ensure the text contains 'Project Status – From To'."
        )

    presentation.save(str(output_path))
    return output_path


def _fill_section_for_team(
    presentation: Presentation,
    team_name: str,
    items: list[dict],
    start_slide_idx: int,
    start_counter: int,
    from_text: str,
    to_text: str,
    upcoming_release: str | None,
) -> tuple[int, int, int]:
    if not items:
        return 0, start_counter, start_slide_idx

    apply_team_text_to_slide(presentation.slides[start_slide_idx], from_text, to_text, team_name)

    filled_total = 0
    insert_position = start_slide_idx
    remaining = items
    counter = start_counter

    while remaining:
        filled, counter = fill_user_story_table(presentation, remaining, insert_position, counter, upcoming_release)
        filled_total += filled
        remaining = remaining[filled:]
        if remaining:
            new_slide_idx = clone_slide(presentation, 1, insert_position)
            apply_team_text_to_slide(presentation.slides[new_slide_idx], from_text, to_text, team_name)
            insert_position = new_slide_idx

    return filled_total, counter, insert_position


def fill_weekly_status_ppt(
    template_path: Path,
    output_path: Path,
    ref_date: date | None = None,
    date_format: str = "%b %d, %Y",
    team: str | None = None,
    progress_callback: Callable[[str], None] | None = None,
) -> tuple[Path, int, int]:
    _emit_progress(progress_callback, "Task: Starting weekly PPT generation.")
    monday, friday = week_range(ref_date)
    from_text, to_text = format_range(monday, friday, date_format)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    _emit_progress(progress_callback, "Task: Loading PPT template. (In Progress)")
    presentation = Presentation(str(template_path))

    if len(presentation.slides) > 0:
        apply_team_text_to_slide(presentation.slides[0], from_text, to_text, team or "")
        _emit_progress(progress_callback, "Task: Updated cover slide date and team details. (Completed)")

    upcoming_release = get_milestone_release_mapping(ref_date)
    total_filled = 0
    _emit_progress(progress_callback, "Task: Resolved milestone mapping for this week. (Completed)")

    if team is None:
        debug("Fetching Eagle User Stories...")
        _emit_progress(progress_callback, "Task: Fetching Eagle user stories from Rally. (In Progress)")
        eagle_us = get_user_stories_from_rally_api(iteration="current", team="Eagle", ai=True, milestone=True)
        filled, _, insert_position = _fill_section_for_team(
            presentation,
            team_name="Eagle",
            items=eagle_us,
            start_slide_idx=1,
            start_counter=1,
            from_text=from_text,
            to_text=to_text,
            upcoming_release=upcoming_release,
        )
        total_filled += filled
        _emit_progress(progress_callback, f"Task: Filled Eagle section with {filled} item(s). (Completed)")

        debug("Fetching Panther Features...")
        _emit_progress(progress_callback, "Task: Fetching Panther feature items from Rally. (In Progress)")
        panther_features = get_feature_items(team="Panther", iteration="current", ai=True, milestone=True)
        panther_slide_idx = clone_slide(presentation, 1, insert_position)
        filled_panther, _, _ = _fill_section_for_team(
            presentation,
            team_name="Panther",
            items=panther_features,
            start_slide_idx=panther_slide_idx,
            start_counter=1,
            from_text=from_text,
            to_text=to_text,
            upcoming_release=upcoming_release,
        )
        total_filled += filled_panther
        _emit_progress(progress_callback, f"Task: Filled Panther section with {filled_panther} item(s). (Completed)")

    elif team.upper() == "PANTHER":
        _emit_progress(progress_callback, f"Task: Fetching Panther feature items from Rally. (In Progress)")
        items = get_feature_items(team="Panther", iteration="current", ai=True, milestone=True)
        for slide in presentation.slides:
            apply_team_text_to_slide(slide, from_text, to_text, "Panther")

        filled, _, _ = _fill_section_for_team(
            presentation,
            team_name="Panther",
            items=items,
            start_slide_idx=1,
            start_counter=1,
            from_text=from_text,
            to_text=to_text,
            upcoming_release=upcoming_release,
        )
        total_filled += filled
        _emit_progress(progress_callback, f"Task: Filled Panther section with {filled} item(s). (Completed)")

    else:
        _emit_progress(progress_callback, f"Task: Fetching {team} user stories from Rally. (In Progress)")
        items = get_user_stories_from_rally_api(iteration="current", team=team, ai=True, milestone=True)
        for slide in presentation.slides:
            apply_team_text_to_slide(slide, from_text, to_text, team)

        filled, _, _ = _fill_section_for_team(
            presentation,
            team_name=team,
            items=items,
            start_slide_idx=1,
            start_counter=1,
            from_text=from_text,
            to_text=to_text,
            upcoming_release=upcoming_release,
        )
        total_filled += filled
        _emit_progress(progress_callback, f"Task: Filled {team} section with {filled} item(s). (Completed)")

    pages_created = len(presentation.slides) - len(Presentation(str(template_path)).slides) + 1
    _emit_progress(progress_callback, "Task: Saving generated PPT file. (In Progress)")
    presentation.save(str(output_path))
    _emit_progress(progress_callback, "Task: PPT generation completed. (Done)")
    return output_path, total_filled, pages_created

